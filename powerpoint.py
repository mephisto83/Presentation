import json
from pptx import Presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.text.text import Font
from pptx.dml.color import _NoneColor
from pptx.chart.chart import Chart
from pptx.chart.series import BarSeries, LineSeries

from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Length, Centipoints, Cm, Emu, Inches, Mm, Pt, Px
import pptx
import inspect
prs = Presentation('existing.pptx')
slide = prs.slides[0]
print("shape count : {}".format(len(slide.shapes)))
output_slides = []
def iter_rPrs(txBody):
    for p in txBody.p_lst:
        for elm in p.content_children:
            yield elm.get_or_add_rPr()
        # generate a:endParaRPr for each <a:p> element
        yield p.get_or_add_endParaRPr()
def captureColor(color):
    result = {}
    print(color)
    print(color.type)
    print(color._color)
    if not isinstance(color._color, _NoneColor):
        result["rgb"] = color.rgb
    return result    
def captureFill(fill):
    result = {}
    result["fore_color"] = captureColor(fill.fore_color)
    return result
for shape in slide.shapes:
    panel = {}
    output_slides.append(panel)
    if isinstance(shape, GraphicFrame):
        print("is a graphic frame ------------------------------------------------------")
        if shape.has_chart:
            print("has a chart")
            chart_part = shape.chart_part
            print(chart_part)
            chart = chart_part.chart;
            print(chart)
            print("chart type {}".format(chart.chart_type))
            if(chart.has_legend):
                print("chart has a legend")
                legend = chart.legend
                print(legend.font.name)
                legendInfo = {"font": legend.font.name}
                panel["legend"] = legendInfo
                print(legend)
            panel["chart_style"] = chart.chart_style   
            print(chart.plots)
            print(chart.series)
            print(len(chart.series))
            c_series = []
            panel["series"] = c_series
            for serie in chart.series:
                print(serie)
                c_serie = {}
                c_series.append(c_serie)
                
                c_serie["smooth"] = False
                if isinstance(serie, LineSeries):
                    print("Smooth : {}".format(serie.smooth))
                    c_serie["smooth"] = serie.smooth
                    
                c_serie["fill"] = None
                if isinstance(serie, BarSeries):
                    print("Fill : {}".format(serie.fill))
                    print("Fill for color : {}".format(serie.fill._fill))
                    fill_format = {}
                    c_serie["fill"] = fill_format
                    if serie.line != None:
                        fill_format["color"] = captureColor(serie.line.color)
                        fill_format["fill"] = captureFill(serie.line.fill)
                        fill_format["width"] = serie.line.width
                        
                    # c_serie["fill"] = serie.fill
                c_serie["values"] = serie.values
                print(serie.values)
                c_serie["name"] = serie.name
                print(serie.name);
                c_serie["index"] = serie.index
                print(serie.index)
                  
        if shape.has_table:
            print("has a table")
    if isinstance(shape, Picture):
        panel["filename"] = shape.image.filename
        panel["fileext"] = shape.image.ext
        panel["contenttype"] = shape.image.content_type
        # print(shape.image.blob)
        with open(shape.name + "." + shape.image.ext, "wb") as blob_file:
            blob_file.write(shape.image.blob)
        print(shape.image)
    if hasattr(shape, "text_frame"):
        fff = Font(shape.text_frame._element.p_lst[0].content_children[0].get_or_add_rPr())
        print(fff)
        print("name {} ".format(fff.name))
        panel["font"] = fff.name
        if fff.size != None:
            print("size {} ".format(Length(fff.size).pt))
            panel["size"] = Length(fff.size).pt
        print("fill {} ".format(fff.fill))
        for txBody in shape._element:
            try:
                if hasattr(txBody, "p_lst"):
                    print("txBody")
                    print(txBody)
                    for rPr in iter_rPrs(txBody):
                        print("font")
                        fon = Font(rPr)
                        
                
            except:
                print("error")
                continue
        panel["marginBottom"] = Length(shape.text_frame.margin_bottom).pt
        panel["marginTop"] = Length(shape.text_frame.margin_top).pt
        panel["marginLeft"] = Length(shape.text_frame.margin_left).pt
        panel["marginRight"] = Length(shape.text_frame.margin_right).pt
    if hasattr(shape, "text"):
        print("text {}".format(shape.text))
        panel["text"] = shape.text
    if hasattr(shape, "height"):
        print("height {}".format(Length(shape.height).pt))
        panel["height"] = Length(shape.height).pt
    if hasattr(shape, "width"):
        print("width {}".format(Length(shape.width).pt))
        panel["width"] = Length(shape.width).pt
    if hasattr(shape, "left"):
        print("left {}".format(Length(shape.left).pt))
        panel["left"] = Length(shape.left).pt
    if hasattr(shape, "top"):
        print("top {}".format(Length(shape.top).pt))
        panel["top"] = Length(shape.top).pt
    if hasattr(shape, "rotation"):
        print("rotation {}".format(shape.rotation))
        panel["rotation"] = shape.rotation
    if hasattr(shape, "name"):
        print("name {}".format(shape.name))
        panel["name"] = shape.name

print(json.dumps(output_slides, sort_keys=True, indent=4, separators=(',', ': ')))
# prs.save('test.pptx')
with open("Output.json", "w") as text_file:
    text_file.write(json.dumps(output_slides, sort_keys=True, indent=4, separators=(',', ': ')))
